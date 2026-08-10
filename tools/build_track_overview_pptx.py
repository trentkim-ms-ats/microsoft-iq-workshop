from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pathlib import Path

BASE = Path(__file__).resolve().parent.parent / "common" / "docs"


def add_slide(prs, title, one_liner, flow_items, policy_rows, tech_points, conclusion, palette):
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Background
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(*palette['bg'])

    # Title
    tx = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.2), Inches(0.7))
    tf = tx.text_frame
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = title
    r.font.size = Pt(30)
    r.font.bold = True
    r.font.color.rgb = RGBColor(*palette['title'])

    # One-liner
    tx2 = slide.shapes.add_textbox(Inches(0.55), Inches(1.0), Inches(12.0), Inches(0.6))
    tf2 = tx2.text_frame
    p2 = tf2.paragraphs[0]
    r2 = p2.add_run()
    r2.text = one_liner
    r2.font.size = Pt(16)
    r2.font.color.rgb = RGBColor(*palette['text'])

    # Flow box
    flow_box = slide.shapes.add_shape(1, Inches(0.5), Inches(1.7), Inches(12.3), Inches(1.1))
    flow_box.fill.solid()
    flow_box.fill.fore_color.rgb = RGBColor(*palette['panel'])
    flow_box.line.color.rgb = RGBColor(*palette['accent'])
    ftxt = flow_box.text_frame
    ftxt.text = 'Architecture Flow'
    ftxt.paragraphs[0].runs[0].font.bold = True
    ftxt.paragraphs[0].runs[0].font.size = Pt(12)
    ftxt.paragraphs[0].runs[0].font.color.rgb = RGBColor(*palette['title'])
    p = ftxt.add_paragraph()
    p.text = ' -> '.join(flow_items)
    p.level = 0
    p.runs[0].font.size = Pt(13)
    p.runs[0].font.color.rgb = RGBColor(*palette['text'])

    # Policy matrix
    matrix = slide.shapes.add_shape(1, Inches(0.5), Inches(3.0), Inches(6.0), Inches(3.7))
    matrix.fill.solid()
    matrix.fill.fore_color.rgb = RGBColor(*palette['panel'])
    matrix.line.color.rgb = RGBColor(*palette['accent'])
    mt = matrix.text_frame
    mt.text = 'Policy Matrix'
    mt.paragraphs[0].runs[0].font.bold = True
    mt.paragraphs[0].runs[0].font.size = Pt(12)
    mt.paragraphs[0].runs[0].font.color.rgb = RGBColor(*palette['title'])
    for row in policy_rows:
        p = mt.add_paragraph()
        p.text = f"- {row}"
        p.runs[0].font.size = Pt(11)
        p.runs[0].font.color.rgb = RGBColor(*palette['text'])

    # Tech points
    tech = slide.shapes.add_shape(1, Inches(6.8), Inches(3.0), Inches(6.0), Inches(2.4))
    tech.fill.solid()
    tech.fill.fore_color.rgb = RGBColor(*palette['panel'])
    tech.line.color.rgb = RGBColor(*palette['accent'])
    tt = tech.text_frame
    tt.text = 'Implementation Points'
    tt.paragraphs[0].runs[0].font.bold = True
    tt.paragraphs[0].runs[0].font.size = Pt(12)
    tt.paragraphs[0].runs[0].font.color.rgb = RGBColor(*palette['title'])
    for t in tech_points:
        p = tt.add_paragraph()
        p.text = f"- {t}"
        p.runs[0].font.size = Pt(11)
        p.runs[0].font.color.rgb = RGBColor(*palette['text'])

    # Conclusion
    concl = slide.shapes.add_shape(1, Inches(6.8), Inches(5.6), Inches(6.0), Inches(1.1))
    concl.fill.solid()
    concl.fill.fore_color.rgb = RGBColor(*palette['accent_light'])
    concl.line.color.rgb = RGBColor(*palette['accent'])
    ct = concl.text_frame
    ct.text = 'Executive Conclusion'
    ct.paragraphs[0].runs[0].font.bold = True
    ct.paragraphs[0].runs[0].font.size = Pt(12)
    ct.paragraphs[0].runs[0].font.color.rgb = RGBColor(*palette['title'])
    p = ct.add_paragraph()
    p.text = conclusion
    p.runs[0].font.size = Pt(11)
    p.runs[0].font.color.rgb = RGBColor(*palette['text'])


def make_track1():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    add_slide(
        prs,
        'Track1 Overview | FabricIQ Foundation with Ontology',
        'Track1 establishes trusted structured data and ontology paths required for Track2 quality gates, Track3 WebIQ scope, and Track4 FoundryIQ grounded responses.',
        ['Source Data', 'Profiling', 'Standardization', 'Ontology Modeling', 'Mapping', 'Validation', 'Handover'],
        [
            'Normal: standardized schema + ontology mapping completed (READY)',
            'Schema mismatch: transform rules applied (PARTIAL)',
            'Reference integrity failure: error rows isolated (PARTIAL)',
            'Critical key duplication: pipeline blocked for correction (BLOCKED)',
        ],
        [
            'Reusable SQL profiling templates',
            'Fixed key/type/code standards',
            'Entity/Relation cardinality controls',
            'Source-to-ontology lineage mapping artifacts',
        ],
        'Track1 outcome is not raw data cleanup; it is a reliable FabricIQ-ready semantic foundation.',
        {
            'bg': (248, 251, 255),
            'panel': (235, 244, 255),
            'accent': (0, 102, 204),
            'accent_light': (220, 236, 255),
            'title': (0, 51, 102),
            'text': (23, 37, 58),
        },
    )
    prs.save(str(BASE / "Track1_Overview.pptx"))


def make_track2():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    add_slide(
        prs,
        'Track2 Overview | WorkIQ Quality Gate and Deployment Readiness',
        'Track2 validates business-evidence data quality and reproducible M365 deployment for WorkIQ retrieval in Track4 FoundryIQ.',
        ['Track1 Handover', 'Sample Generation', 'M365 Deployment', '8-Item Quality Gate', 'Scoring', 'Manifest', 'Track4 Input'],
        [
            'Normal: generation/deployment/manifest complete (PASS)',
            'Partial channel failure: retry with channel status tracking (PARTIAL)',
            'Score < 75: regenerate and re-validate failed items (PARTIAL)',
            'ACL failure or manifest mismatch: handover blocked (BLOCKED)',
        ],
        [
            'Seed+Extended distribution consistency controls',
            'One-click generate/execute deployment path',
            '8-item gate with 6-pass minimum rule',
            'Manifest-based evidence traceability for Track4 FoundryIQ',
        ],
        'Track2 outcome is not sample volume; it is policy-compliant WorkIQ evidence readiness with reproducibility.',
        {
            'bg': (248, 255, 249),
            'panel': (234, 248, 236),
            'accent': (20, 128, 61),
            'accent_light': (217, 242, 223),
            'title': (14, 92, 43),
            'text': (22, 52, 36),
        },
    )
    prs.save(str(BASE / "Track2_Overview.pptx"))


def make_track4_foundry():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    add_slide(
        prs,
        'Track4 Overview | FoundryIQ Operational Agent Validation',
        'Track4 verifies routed, evidence-grounded responses under normal and failure conditions with explicit fallback policies.',
        ['Question', 'FoundryIQ Orchestrator', 'FabricIQ Tool', 'WorkIQ Tool', 'Merge Response', 'Quality Gate', 'Ops Report'],
        [
            'Normal: Tool A/B merged response (PASS)',
            'Tool A fail: evidence-only limited response (PARTIAL)',
            'Tool B fail: metric-only limited response (PARTIAL)',
            'Both fail: block output with recovery guidance (BLOCKED)',
            'Transient errors: initial call + 5s/10s/20s retries (max 4 attempts)',
        ],
        [
            'Simulation/live dual-mode execution',
            'Adapter contract-based tool invocation',
            'Pass/partial/blocked response grading',
            'Run context and retry audit trace storage',
        ],
        'Track4 outcome is not a flashy demo; it is a testable and resilient operational agent system.',
        {
            'bg': (255, 250, 245),
            'panel': (255, 240, 225),
            'accent': (194, 87, 0),
            'accent_light': (255, 230, 204),
            'title': (128, 58, 0),
            'text': (62, 35, 12),
        },
    )
    prs.save(str(BASE / "Track4_Overview.pptx"))


if __name__ == '__main__':
    make_track1()
    make_track2()
    make_track4_foundry()
    print('created')
