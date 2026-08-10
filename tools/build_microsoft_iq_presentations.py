"""Build the current Microsoft IQ architecture presentations."""

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


OUTPUT_DIR = Path(__file__).resolve().parent.parent / "common" / "docs"
FULL_DECK = OUTPUT_DIR / "Microsoft_IQ_Workshop_Full_Architecture.pptx"
DETAILED_DECK = OUTPUT_DIR / "Microsoft_IQ_Workshop_Architecture_Detailed.pptx"

SLIDE_W = 13.333
SLIDE_H = 7.5

NAVY = "111827"
INK = "172033"
MUTED = "5D6778"
LIGHT = "F6F7FB"
WHITE = "FFFFFF"
BLUE = "2563EB"
BLUE_LIGHT = "EAF1FF"
GREEN = "15966A"
GREEN_LIGHT = "E9F7F1"
PURPLE = "7C3AED"
PURPLE_LIGHT = "F1EBFF"
ORANGE = "D97706"
ORANGE_LIGHT = "FFF1E3"
RED = "BE3455"
RED_LIGHT = "FCECEF"
GOLD = "F5C451"

TRACK_COLORS = [
    (BLUE, BLUE_LIGHT),
    (GREEN, GREEN_LIGHT),
    (PURPLE, PURPLE_LIGHT),
    (ORANGE, ORANGE_LIGHT),
]


def rgb(value: str) -> RGBColor:
    return RGBColor.from_string(value)


def new_deck() -> Presentation:
    prs = Presentation()
    prs.slide_width = Inches(SLIDE_W)
    prs.slide_height = Inches(SLIDE_H)
    return prs


def add_text(
    slide,
    text: str,
    x: float,
    y: float,
    w: float,
    h: float,
    *,
    size: float = 14,
    color: str = INK,
    bold: bool = False,
    font: str = "Aptos",
    align=PP_ALIGN.LEFT,
    valign=MSO_ANCHOR.TOP,
    margin: float = 0,
):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.margin_left = Inches(margin)
    frame.margin_right = Inches(margin)
    frame.margin_top = Inches(margin)
    frame.margin_bottom = Inches(margin)
    frame.vertical_anchor = valign
    paragraph = frame.paragraphs[0]
    paragraph.alignment = align
    paragraph.space_after = Pt(0)
    run = paragraph.add_run()
    run.text = text
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = rgb(color)
    return box


def add_rect(
    slide,
    x: float,
    y: float,
    w: float,
    h: float,
    *,
    fill: str = WHITE,
    line: str = "D8DDE7",
    radius: bool = True,
    line_width: float = 1,
):
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
    shape = slide.shapes.add_shape(
        shape_type, Inches(x), Inches(y), Inches(w), Inches(h)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb(fill)
    shape.line.color.rgb = rgb(line)
    shape.line.width = Pt(line_width)
    return shape


def add_card(
    slide,
    *,
    x: float,
    y: float,
    w: float,
    h: float,
    title: str,
    body: str,
    accent: str,
    fill: str = WHITE,
    label: str | None = None,
    title_size: float = 16,
    body_size: float = 12.5,
):
    add_rect(slide, x, y, w, h, fill=fill, line=accent, line_width=1.1)
    if label:
        badge_w = min(max(0.48, 0.12 * len(label)), w - 0.4)
        add_rect(
            slide,
            x + 0.18,
            y + 0.18,
            badge_w,
            0.28,
            fill=accent,
            line=accent,
            radius=True,
        )
        add_text(
            slide,
            label,
            x + 0.18,
            y + 0.215,
            badge_w,
            0.18,
            size=9,
            color=WHITE,
            bold=True,
            align=PP_ALIGN.CENTER,
        )
        title_y = y + 0.58
    else:
        title_y = y + 0.25
    add_text(
        slide,
        title,
        x + 0.22,
        title_y,
        w - 0.44,
        0.42,
        size=title_size,
        bold=True,
    )
    add_text(
        slide,
        body,
        x + 0.22,
        title_y + 0.55,
        w - 0.44,
        h - (title_y - y) - 0.72,
        size=body_size,
        color=MUTED,
    )


def add_banner(
    slide,
    text: str,
    *,
    x: float = 0.75,
    y: float = 6.22,
    w: float = 11.83,
    h: float = 0.47,
    fill: str = NAVY,
    color: str = WHITE,
    size: float = 11.5,
):
    add_rect(slide, x, y, w, h, fill=fill, line=fill, radius=True)
    add_text(
        slide,
        text,
        x + 0.18,
        y + 0.03,
        w - 0.36,
        h - 0.06,
        size=size,
        color=color,
        bold=True,
        align=PP_ALIGN.CENTER,
        valign=MSO_ANCHOR.MIDDLE,
    )


def add_arrow(slide, x: float, y: float, w: float = 0.24, h: float = 0.24, color: str = BLUE):
    arrow = slide.shapes.add_shape(
        MSO_SHAPE.RIGHT_ARROW, Inches(x), Inches(y), Inches(w), Inches(h)
    )
    arrow.fill.solid()
    arrow.fill.fore_color.rgb = rgb(color)
    arrow.line.fill.background()


def add_header(
    slide,
    *,
    kicker: str,
    title: str,
    subtitle: str,
    number: int,
    source: str,
):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = rgb(LIGHT)
    add_text(slide, kicker, 0.58, 0.28, 4.8, 0.24, size=10, color=ORANGE, bold=True)
    add_text(
        slide,
        str(number),
        12.2,
        0.24,
        0.5,
        0.3,
        size=11,
        color=MUTED,
        bold=True,
        align=PP_ALIGN.RIGHT,
    )
    add_text(slide, title, 0.55, 0.64, 12.0, 0.5, size=27, bold=True)
    add_text(slide, subtitle, 0.58, 1.22, 12.0, 0.44, size=13.5, color=MUTED)
    add_text(slide, f"Source: {source}", 0.58, 7.08, 10.8, 0.18, size=8.5, color="8A93A3")


def add_status_card(
    slide,
    *,
    x: float,
    y: float,
    title: str,
    status: str,
    body: str,
    color: str,
    fill: str,
):
    add_card(
        slide,
        x=x,
        y=y,
        w=3.86,
        h=1.7,
        title=title,
        body=body,
        accent=color,
        fill=fill,
        label=status,
        title_size=15,
        body_size=11.8,
    )


def build_full_deck() -> None:
    prs = new_deck()

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = rgb(NAVY)
    add_text(
        slide,
        "ARCHITECTURE EDITION · JULY 2026",
        0.66,
        0.48,
        5.6,
        0.24,
        size=10,
        color=GOLD,
        bold=True,
    )
    add_text(
        slide,
        "Microsoft IQ Workshop 전체 아키텍처",
        0.64,
        1.04,
        11.8,
        0.7,
        size=31,
        color=WHITE,
        bold=True,
    )
    add_text(
        slide,
        "FabricIQ 정형 신뢰성, WorkIQ 권한 기반 업무 근거, WebIQ 공개 웹 citation,\nFoundryIQ 결합·평가를 하나의 실행 가능한 워크숍으로 연결합니다.",
        0.66,
        1.91,
        11.5,
        0.75,
        size=15,
        color="C9D1E2",
    )
    stats = [
        ("600분", "권장 2일 학습 여정"),
        ("4 Tracks", "Fabric · Work · Web · Foundry"),
        ("3 Handoffs", "책임별 입력 계약"),
        ("1 Briefing", "수치 + 내부/외부 근거 + 조치"),
    ]
    for index, (value, label) in enumerate(stats):
        x = 0.66 + index * 3.04
        add_rect(slide, x, 3.35, 2.72, 1.05, fill="1B2842", line="2B3B5B")
        add_text(slide, value, x + 0.2, 3.54, 2.32, 0.3, size=18, color=WHITE, bold=True)
        add_text(slide, label, x + 0.2, 3.94, 2.32, 0.22, size=10.5, color="BFC9DA")
    add_banner(
        slide,
        "Ontology = 네 구성요소가 같은 비즈니스 대상을 가리키는 공통 어휘",
        x=0.66,
        y=5.35,
        w=12.0,
        h=0.58,
        fill=WHITE,
        color=INK,
        size=13,
    )
    add_text(slide, "TRACK OUTCOME", 0.72, 5.56, 1.5, 0.14, size=8.5, color=ORANGE, bold=True)
    add_text(slide, "Microsoft IQ Workshop", 0.66, 7.05, 3.0, 0.2, size=8.5, color="7C879B")
    add_text(slide, "1", 12.2, 7.02, 0.4, 0.2, size=9, color="7C879B", align=PP_ALIGN.RIGHT)

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(
        slide,
        kicker="SYSTEM CONTEXT",
        title="네 구성요소 Microsoft IQ 아키텍처",
        subtitle="FoundryIQ가 질문을 라우팅하고 FabricIQ, WorkIQ, WebIQ가 책임별 근거를 제공합니다.",
        number=2,
        source="AGENTS.md · integrated plan",
    )
    context_cards = [
        ("FabricIQ", "Lakehouse 정형 데이터\nOntology 관계 경로\n검증 가능한 KPI"),
        ("WorkIQ", "SharePoint · OneDrive\nOutlook · Teams\nACL 적용 내부 근거"),
        ("WebIQ", "공개 웹 citation\n관찰 시각 · scope\nfactStatus · limitations"),
        ("FoundryIQ", "Tool A/B/C 라우팅\n정책 결합 · 평가\n브리핑 + 조치"),
    ]
    for i, (title, body) in enumerate(context_cards):
        x = 0.58 + i * 3.04
        accent, fill = TRACK_COLORS[i]
        add_card(
            slide,
            x=x,
            y=2.22,
            w=2.82,
            h=3.45,
            title=title,
            body=body,
            accent=accent,
            fill=WHITE,
            label=str(i + 1),
            body_size=13,
        )
        if i < 3:
            add_arrow(slide, x + 2.86, 3.7, color=accent)
    add_banner(
        slide,
        "공통 semantic keys · Campaign · Product · Customer · Order · Payment · Shipment · Return",
        y=6.18,
    )

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(
        slide,
        kicker="RESPONSIBILITY CONTRACT",
        title="책임과 금지 경계를 분리",
        subtitle="근거 없는 대체를 허용하지 않는 것이 Microsoft IQ 품질 계약의 핵심입니다.",
        number=3,
        source="AGENTS.md · track1–4 workbooks",
    )
    contract_cards = [
        ("FabricIQ", "책임 · 정형 KPI와 Ontology\n\n금지 · 문서에서 수치 역산\n\n증명 · 기준값 · lineage"),
        ("WorkIQ", "책임 · ACL 적용 내부 검색\n\n금지 · 공개 웹을 내부 근거로 사용\n\n증명 · original link · coverage"),
        ("WebIQ", "책임 · 공개 URL citation\n\n금지 · 출처 없는 최신 사실·내부 원인 단정\n\n증명 · scope · limitations"),
        ("FoundryIQ", "책임 · 라우팅·결합·평가\n\n금지 · 숫자·링크·인과 생성\n\n증명 · sourceTrace · warnings"),
    ]
    for i, (title, body) in enumerate(contract_cards):
        x = 0.58 + i * 3.04
        accent, fill = TRACK_COLORS[i]
        add_card(
            slide,
            x=x,
            y=2.0,
            w=2.82,
            h=4.65,
            title=title,
            body=body,
            accent=accent,
            fill=fill,
            body_size=12.3,
        )

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(
        slide,
        kicker="TRACK 1 · FABRICIQ",
        title="정형 데이터에서 AI 친화적 Ontology까지",
        subtitle="참가자는 P1 탐지 쿼리 대신 데이터 구조·품질 개념·표준화 영향·의미 경로를 이해합니다.",
        number=4,
        source="track1/WORKBOOK.md",
    )
    stages = [
        ("질문", "Q1–Q5\n질문–테이블 연결"),
        ("Orientation", "14개 테이블\n품질 개념과 영향"),
        ("표준화", "키 · 타입 · 코드\n날짜 · 통화"),
        ("Ontology", "14 엔터티\n20 관계"),
        ("Mapping", "3단 매핑\n의미 경로"),
    ]
    for i, (title, body) in enumerate(stages):
        x = 0.58 + i * 2.42
        add_card(
            slide,
            x=x,
            y=2.25,
            w=2.18,
            h=3.25,
            title=title,
            body=body,
            accent=BLUE,
            fill=WHITE,
            label=str(i + 1),
            title_size=14.2,
            body_size=12.2,
        )
        if i < 4:
            add_arrow(slide, x + 2.21, 3.7, w=0.18, h=0.2, color=BLUE)
    add_banner(
        slide,
        "TRACK2 HANDOFF · Workspace/Ontology IDs · 핵심 경로 · Entity 매핑 · 검색 키 · 구현 제한",
        y=6.1,
        size=11,
    )

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(
        slide,
        kicker="TRACK 2 · WORKIQ",
        title="M365 업무 근거와 크로스 소스 품질 게이트",
        subtitle="같은 Ontology 키로 4대 M365 소스를 찾고 ACL·원본 링크·품질 점수를 FoundryIQ에 전달합니다.",
        number=5,
        source="track2/WORKBOOK.md",
    )
    work_cards = [
        ("SharePoint · 15", "기획서 · 리포트 · 플레이북", GREEN, WHITE),
        ("Outlook · 15", "리더십 메일 · 이슈 · 의사결정", GREEN, WHITE),
        ("Teams · 18 / 55", "CS · 재고 · 물류 대화", GREEN, WHITE),
        ("OneDrive · 12", "회의록 · 브리핑 · 정책", GREEN, WHITE),
        ("QUALITY GATE", "정확성 · 완전성 · 일관성 · 유효성\n중복성 · 참조무결성 · 적시성 · 추적성", BLUE, BLUE_LIGHT),
        ("FOUNDRY INPUT", "인덱스 · 품질 점수\n근거 링크 5건 · 재현 질의 3개\nlive ACL 증적", ORANGE, ORANGE_LIGHT),
    ]
    for i, (title, body, accent, fill) in enumerate(work_cards):
        col, row = i % 3, i // 3
        add_card(
            slide,
            x=0.58 + col * 4.08,
            y=2.0 + row * 2.28,
            w=3.82,
            h=2.0,
            title=title,
            body=body,
            accent=accent,
            fill=fill,
            title_size=14.5,
            body_size=11.8,
        )

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(
        slide,
        kicker="TRACK 3 · WEBIQ",
        title="공개 웹 최신 근거를 citation 계약으로 추가",
        subtitle="기존 시나리오와 데이터를 바꾸지 않고 외부에서 확인할 수 있는 질문만 안전하게 분리합니다.",
        number=6,
        source="track3/WORKBOOK.md",
    )
    web_steps = [
        ("Safe query", "민감 정보 제거\n공개 확인 범위", "10분"),
        ("Search", "Web Search 또는 fixture\ncitation ≥2", "15분"),
        ("Quality", "권위 · 현재성\nscope · factStatus\nlimitations", "10분"),
        ("Handoff", "TRACK4_FOUNDRYIQ_HANDOFF_PACKAGE", "10분"),
    ]
    for i, (title, body, label) in enumerate(web_steps):
        x = 0.58 + i * 3.04
        add_card(
            slide,
            x=x,
            y=2.18,
            w=2.82,
            h=3.65,
            title=title,
            body=body,
            accent=PURPLE,
            fill=WHITE,
            label=label,
            body_size=12.8,
        )
        if i < 3:
            add_arrow(slide, x + 2.86, 3.85, color=PURPLE)
    add_banner(
        slide,
        "WebIQ는 공개 사실을 제공하며 WorkIQ 내부 근거나 FabricIQ 정형 수치를 대체하지 않음",
        y=6.2,
        size=11,
    )

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(
        slide,
        kicker="TRACK 4 · FOUNDRYIQ",
        title="런타임 오케스트레이션과 응답 계약",
        subtitle="세 도구를 호출하고 소스 책임을 보존한 채 최종 문장, 경고, 조치안을 만듭니다.",
        number=7,
        source="track4/WORKBOOK.md · track4/data/README.md",
    )
    orchestration = [
        ("QUESTION", "scenarioId · question\nsemanticKeys[]", ORANGE),
        ("TOOL A", "FabricIQ\nmetrics · sourceTrace", BLUE),
        ("TOOL B", "WorkIQ\nlinks · ACL · coverage", GREEN),
        ("TOOL C", "WebIQ\ncitations · scope · limits", PURPLE),
        ("MERGE", "근거 책임 보존\nwarnings · fallback", ORANGE),
        ("EVALUATE", "PASS · PARTIAL · BLOCKED\nstrict report", RED),
    ]
    for i, (title, body, accent) in enumerate(orchestration):
        col, row = i % 3, i // 3
        add_card(
            slide,
            x=0.58 + col * 4.08,
            y=2.0 + row * 2.32,
            w=3.82,
            h=2.02,
            title=title,
            body=body,
            accent=accent,
            fill=WHITE,
            title_size=14.5,
            body_size=12.3,
        )

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(
        slide,
        kicker="EXECUTION PROFILES",
        title="simulation과 live는 같은 정책, 다른 증거 수준",
        subtitle="simulation은 교육·회귀용 reference harness이며 실제 서비스 연결 성공을 의미하지 않습니다.",
        number=8,
        source="AGENTS.md 실행 모드",
    )
    add_card(
        slide,
        x=0.58,
        y=2.0,
        w=5.9,
        h=3.85,
        title="SIMULATION",
        body="오프라인 재현\n\nTrack1 CSV · Track2 manifest\nTrack3 WebIQ fixture\n\n동일 response schema · 동일 strict evaluator\n실제 ACL·네트워크·현재성 미증명",
        accent=PURPLE,
        fill=PURPLE_LIGHT,
        label="교육·회귀",
        title_size=18,
        body_size=13,
    )
    add_card(
        slide,
        x=6.78,
        y=2.0,
        w=5.9,
        h=3.85,
        title="LIVE",
        body="실제 연결 검증\n\nFabricIQ adapter · WorkIQ adapter\nFoundry Web Search\n\n권한·URL·관찰 시각 확인\n오류와 제한을 명시적으로 표면화",
        accent=ORANGE,
        fill=ORANGE_LIGHT,
        label="승인된 검증",
        title_size=18,
        body_size=13,
    )
    add_banner(
        slide,
        "공통 통제 · sourceTrace 분리 · 5s → 10s → 20s · 부분응답 경고 · release version",
        y=6.2,
        size=11.3,
    )

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(
        slide,
        kicker="RESILIENCE & QUALITY",
        title="소스별 fallback 상태 머신",
        subtitle="오류를 숨기지 않고 신뢰 가능한 범위만 제공한 뒤 복구 후 같은 질문으로 재평가합니다.",
        number=9,
        source="AGENTS.md fallback 정책 · track4/WORKBOOK.md",
    )
    status_cards = [
        ("Normal", "PASS", "세 근거 정상\n수치 + 내부 링크 + 외부 citation + 조치", GREEN, GREEN_LIGHT),
        ("fabric-down", "PARTIAL", "정형 수치 미검증\n남은 근거만 제한 제공", ORANGE, ORANGE_LIGHT),
        ("work-down", "PARTIAL", "내부 업무 근거 없음\n수치 + 공개 근거", ORANGE, ORANGE_LIGHT),
        ("web-down", "PARTIAL", "현재 외부 확인 없음\n수치 + 내부 근거", ORANGE, ORANGE_LIGHT),
        ("internal-down", "BLOCKED", "웹만으로 내부 원인 분석 금지\n최종 답변 중단", RED, RED_LIGHT),
        ("all-down", "BLOCKED", "복구 조치만 제공", RED, RED_LIGHT),
    ]
    for i, (title, status, body, color, fill) in enumerate(status_cards):
        col, row = i % 3, i // 3
        add_status_card(
            slide,
            x=0.58 + col * 4.08,
            y=2.0 + row * 2.02,
            title=title,
            status=status,
            body=body,
            color=color,
            fill=fill,
        )

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(
        slide,
        kicker="WORKSHOP JOURNEY",
        title="600분 권장 학습 여정",
        subtitle="입문자는 기존 480분 기반을 먼저 경험한 뒤 WebIQ와 네 구성요소 통합을 추가합니다.",
        number=10,
        source="Microsoft_IQ_Workshop_Integrated_Plan.md",
    )
    add_card(
        slide,
        x=0.58,
        y=2.0,
        w=5.9,
        h=3.9,
        title="DAY 1 · 480분",
        body="원본 Microsoft IQ 베이스라인\n\n오프닝 · 개념\nTrack1 FabricIQ · 150분\nTrack2 WorkIQ · 110분\n기존 FoundryIQ 심화 경로\n통합 프로젝트 · 리뷰",
        accent=BLUE,
        fill=BLUE_LIGHT,
        title_size=17,
        body_size=13,
    )
    add_card(
        slide,
        x=6.78,
        y=2.0,
        w=5.9,
        h=3.9,
        title="DAY 2 · 120분",
        body="WebIQ-enabled Microsoft IQ\n\n복습 · 책임 분리 · 10분\nTrack3 WebIQ · 45분\nTrack4 네 구성요소 통합 · 45분\n비교 · 발표 · 20분",
        accent=PURPLE,
        fill=PURPLE_LIGHT,
        title_size=17,
        body_size=13,
    )
    add_banner(
        slide,
        "선택 운영 · 1일 압축형 480분 · Track1 150 → Track2 110 → Track3 35 → Track4 60 → 프로젝트 35",
        y=6.22,
        size=10.8,
    )

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(
        slide,
        kicker="OPERATIONS ARCHITECTURE",
        title="재현 가능한 실행·평가·자동화 구조",
        subtitle="데이터, manifest, fixture, 스크립트, evaluator가 같은 계약으로 회귀 가능한 파이프라인을 구성합니다.",
        number=11,
        source="AGENTS.md · track4/data/README.md",
    )
    ops = [
        ("Trigger", "수동 · cron\n승인된 일정"),
        ("Generate", "Q1–Q3\n세 입력 계약"),
        ("Execute", "simulation/live\nTool A/B/C"),
        ("Evaluate", "strict evaluator\n상태·경고"),
        ("Brief", "근거·조치\n리더십 브리핑"),
        ("Notify", "사람 승인 후\nTeams/Outlook"),
    ]
    accents = [BLUE, GREEN, PURPLE, ORANGE, ORANGE, RED]
    for i, (title, body) in enumerate(ops):
        x = 0.58 + i * 2.03
        add_card(
            slide,
            x=x,
            y=2.3,
            w=1.8,
            h=3.2,
            title=title,
            body=body,
            accent=accents[i],
            fill=WHITE,
            label=str(i + 1),
            title_size=13.2,
            body_size=11,
        )
        if i < 5:
            add_arrow(slide, x + 1.83, 3.72, w=0.16, h=0.2, color=accents[i])
    add_banner(
        slide,
        "DELIVERY GATE · strict FAIL → 발송 중단 · PARTIAL → 경고 + 수동 검토",
        y=6.12,
        size=11.3,
    )

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(
        slide,
        kicker="END STATE",
        title="참가자가 마지막에 증명해야 하는 것",
        subtitle="아키텍처 설명에서 끝나지 않고 실행 결과와 근거로 네 구성요소의 책임을 증명합니다.",
        number=12,
        source="track1–4 workbooks · integrated plan",
    )
    outcomes = [
        ("FabricIQ", "이 수치는 어떤 엔터티·관계·속성에서 계산됐는가?\n\n기준값 · semantic path · mapping proof"),
        ("WorkIQ", "이 내부 근거는 누가 볼 수 있고 원문은 어디에 있는가?\n\nACL · source coverage · original link"),
        ("WebIQ", "이 공개 사실은 언제·어디에 적용되며 한계는 무엇인가?\n\nURL · observedAt · scope · limitations"),
        ("FoundryIQ", "세 소스를 어떻게 결합했고 실패 시 무엇을 제한했는가?\n\nsourceTrace · warnings · strict report"),
    ]
    for i, (title, body) in enumerate(outcomes):
        x = 0.58 + i * 3.04
        accent, fill = TRACK_COLORS[i]
        add_card(
            slide,
            x=x,
            y=2.0,
            w=2.82,
            h=4.6,
            title=title,
            body=body,
            accent=accent,
            fill=WHITE,
            title_size=16,
            body_size=12.4,
        )

    prs.save(FULL_DECK)


def build_detailed_deck() -> None:
    prs = new_deck()

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(
        slide,
        kicker="INTEGRATED ARCHITECTURE",
        title="Microsoft IQ Workshop Integrated Architecture",
        subtitle="Track1 FabricIQ → Track2 WorkIQ → Track3 WebIQ → Track4 FoundryIQ",
        number=1,
        source="Microsoft_IQ_Workshop_Integrated_Plan.md",
    )
    add_banner(
        slide,
        "Cross-cutting controls · Security/ACL · Data lineage · Observability · Prompt guardrails · Compliance logging",
        x=0.58,
        y=1.72,
        w=12.17,
        h=0.5,
        size=10.5,
    )
    architecture = [
        ("Source Domain", "ERP/CRM transactions\nCampaign & support events", NAVY),
        ("Track1 · FabricIQ", "Lakehouse data model\nOntology contract\nMapping review + handoff", BLUE),
        ("Track2 · WorkIQ", "M365 indexing\nEvidence links + ACL\nQuality gate: 6/8 PASS", GREEN),
        ("Track3 · WebIQ", "Safe public query\nURL citation + scope\nSource quality + limitations", PURPLE),
        ("Track4 · FoundryIQ", "Tool A/B/C routing\nFallback + strict evaluation\nExecutive response + actions", ORANGE),
    ]
    for i, (title, body, accent) in enumerate(architecture):
        x = 0.58 + i * 2.45
        add_card(
            slide,
            x=x,
            y=2.53,
            w=2.18,
            h=3.3,
            title=title,
            body=body,
            accent=accent,
            fill=WHITE,
            label=str(i),
            title_size=12.5,
            body_size=10.8,
        )
        if i < 4:
            add_arrow(slide, x + 2.21, 3.97, w=0.2, h=0.22, color=accent)
    add_banner(
        slide,
        "Output · Executive brief + traceable internal evidence + cited external evidence + prioritized action backlog",
        x=0.58,
        y=6.16,
        w=12.17,
        h=0.5,
        size=10.7,
    )

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(
        slide,
        kicker="RUNTIME ORCHESTRATION",
        title="Detailed Runtime Orchestration Flow",
        subtitle="From business question intake to source-responsible response generation and audit logging",
        number=2,
        source="AGENTS.md · track4/WORKBOOK.md",
    )
    add_card(
        slide,
        x=0.58,
        y=2.3,
        w=1.62,
        h=1.22,
        title="1 · REQUEST",
        body="Q1/Q2/Q3 input",
        accent=ORANGE,
        fill=WHITE,
        title_size=11.5,
        body_size=10.5,
    )
    add_arrow(slide, 2.28, 2.8, w=0.25, h=0.24, color=ORANGE)
    add_card(
        slide,
        x=2.62,
        y=2.3,
        w=1.72,
        h=1.22,
        title="2 · ROUTER",
        body="Intent classify\nsemantic keys",
        accent=ORANGE,
        fill=ORANGE_LIGHT,
        title_size=11.5,
        body_size=10.3,
    )
    tools = [
        ("TOOL A · FABRICIQ", "Metrics + structured sourceTrace", BLUE, 1.95),
        ("TOOL B · WORKIQ", "ACL evidence + original links", GREEN, 3.43),
        ("TOOL C · WEBIQ", "Public citations + scope + limits", PURPLE, 4.91),
    ]
    for title, body, accent, y in tools:
        add_card(
            slide,
            x=4.72,
            y=y,
            w=3.26,
            h=1.1,
            title=title,
            body=body,
            accent=accent,
            fill=WHITE,
            title_size=11.8,
            body_size=10.5,
        )
    add_arrow(slide, 4.4, 2.8, w=0.25, h=0.24, color=ORANGE)
    add_arrow(slide, 8.08, 3.6, w=0.25, h=0.24, color=ORANGE)
    add_card(
        slide,
        x=8.42,
        y=2.35,
        w=1.8,
        h=3.05,
        title="3 · MERGE",
        body="Preserve source responsibility\n\nAdd warnings and limitations",
        accent=ORANGE,
        fill=ORANGE_LIGHT,
        title_size=12,
        body_size=10.8,
    )
    add_arrow(slide, 10.31, 3.6, w=0.25, h=0.24, color=RED)
    add_card(
        slide,
        x=10.66,
        y=2.35,
        w=2.05,
        h=3.05,
        title="4 · EVALUATE",
        body="PASS\nPARTIAL\nBLOCKED\n\nScore + audit",
        accent=RED,
        fill=RED_LIGHT,
        title_size=12,
        body_size=11,
    )
    add_banner(
        slide,
        "Mandatory runtime gates · evidence validity · citation quality · traceability completeness · policy-compliant warnings · human approval",
        x=0.75,
        y=6.24,
        w=11.83,
        h=0.5,
        size=10.2,
    )

    prs.save(DETAILED_DECK)


if __name__ == "__main__":
    OUTPUT_DIR.mkdir(parents=True, exist_ok=True)
    build_full_deck()
    build_detailed_deck()
    print(f"created: {FULL_DECK}")
    print(f"created: {DETAILED_DECK}")